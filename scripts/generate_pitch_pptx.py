#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""生成群像·星火路演PPT（知乎蓝白风）"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# 知乎蓝白风配色
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ZHIBLUE = RGBColor(0x00, 0x66, 0xFF)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x55, 0x55, 0x55)
LIGHT_GRAY = RGBColor(0x88, 0x88, 0x88)
BG_GRAY = RGBColor(0xF6, 0xF8, 0xFA)
GOLD = RGBColor(0xE2, 0xB0, 0x4A)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide, color=WHITE):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size, bold=False,
                color=DARK, align=PP_ALIGN.CENTER, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox

# ========== 第1页：封面 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
# 左侧几何装饰
shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.5), Inches(-1.5), Inches(5), Inches(5))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xE8, 0xF0, 0xFF)
shape.line.fill.background()
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.2), Inches(2), Inches(2))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xF0, 0xF5, 0xFF)
shape.line.fill.background()

add_textbox(slide, Inches(0), Inches(2.2), Inches(13.333), Inches(1.2),
            '群像·星火', 60, bold=True, color=ZHIBLUE)
add_textbox(slide, Inches(0), Inches(3.5), Inches(13.333), Inches(0.6),
            '让真实发光，让思想变现', 24, color=DARK)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.8), Inches(4.2), Inches(1.7), Inches(0.02))
line.fill.solid()
line.fill.fore_color.rgb = ZHIBLUE
line.line.fill.background()
add_textbox(slide, Inches(0), Inches(4.5), Inches(13.333), Inches(0.5),
            '一个关于创作、连接与被看见的故事', 14, color=LIGHT_GRAY)
# 右下角标签
tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.5), Inches(6.2), Inches(3.0), Inches(0.5))
tag.fill.solid()
tag.fill.fore_color.rgb = WHITE
tag.line.color.rgb = RGBColor(0xCC, 0xDD, 0xFF)
add_textbox(slide, Inches(9.5), Inches(6.25), Inches(3.0), Inches(0.4),
            '知乎黑客松 2026 参赛项目', 12, color=ZHIBLUE, align=PP_ALIGN.CENTER)

# ========== 第2页：我们的初心 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, Inches(0), Inches(0.8), Inches(13.333), Inches(0.9),
            '希望有一天…', 42, bold=True, color=ZHIBLUE)
add_textbox(slide, Inches(0), Inches(1.6), Inches(13.333), Inches(0.4),
            '这是我们做这件事的全部理由', 14, color=LIGHT_GRAY)

hopes = [
    ('人人都是', '创作者'),
    ('我们能成为', '故事的参与者'),
    ('我们能在故事里相遇，', '争论小小细节'),
    ('在玩故事中', '收获快乐和朋友'),
    ('坑了的小说', '被我们填平'),
    ('我们', '都能被看见'),
]
y_start = 2.3
for i, (prefix, highlight) in enumerate(hopes):
    y = y_start + i * 0.65
    add_textbox(slide, Inches(3.5), Inches(y), Inches(3.5), Inches(0.5),
                prefix, 20, color=GRAY, align=PP_ALIGN.RIGHT)
    add_textbox(slide, Inches(7.0), Inches(y), Inches(3.5), Inches(0.5),
                highlight, 20, bold=True, color=ZHIBLUE, align=PP_ALIGN.LEFT)

# ========== 第3页：三个模式 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, Inches(0), Inches(0.8), Inches(13.333), Inches(0.9),
            '三个模式，一个信念', 42, bold=True, color=ZHIBLUE)

cards = [
    ('🤖', '人机模式', '和刘看山聊聊你的脑洞\n一个永远在线的倾听者'),
    ('💬', '双人模式', '和陌生人碰撞出真实火花\n两个真人 + 一个AI催化剂'),
    ('🎭', '故事模式', '选一个角色，揭开历史谜题\n四幕推进，真相由你揭晓'),
]
for i, (icon, title, desc) in enumerate(cards):
    x = 1.5 + i * 4.0
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(x), Inches(2.2), Inches(3.2), Inches(3.2))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(0xCC, 0xDD, 0xFF)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     Inches(x + 1.1), Inches(2.5), Inches(1.0), Inches(1.0))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0xE8, 0xF0, 0xFF)
    circle.line.fill.background()
    add_textbox(slide, Inches(x), Inches(2.65), Inches(3.2), Inches(0.7),
                icon, 28, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(x), Inches(3.6), Inches(3.2), Inches(0.5),
                title, 18, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(x + 0.2), Inches(4.1), Inches(2.8), Inches(1.0),
                desc, 12, color=GRAY, align=PP_ALIGN.CENTER)

add_textbox(slide, Inches(0), Inches(6.3), Inches(13.333), Inches(0.4),
            'AI不是作者，是催化剂。你才是故事的主角。', 13, color=LIGHT_GRAY)

# ========== 第4页：刘看山 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, Inches(0), Inches(0.6), Inches(13.333), Inches(0.9),
            '一个Agent，11种角色', 42, bold=True, color=ZHIBLUE)
add_textbox(slide, Inches(0), Inches(1.4), Inches(13.333), Inches(0.4),
            '刘看山 —— 最懂你的AI副导演', 16, color=GRAY)

roles = [
    'Companion · 陪伴者', 'DM · 主持人', 'Story Fallback · 故事替补',
    'AI副导演 · 剧情把控', 'Catalyst · 催化剂', 'Healer · 疗愈师',
    'Reviewer · 审稿人', 'Summarizer · 总结者', 'Knowledge Feeder · 知识投喂',
    'Mediator · 调解员', 'Creative · 创作助手'
]
for i, role in enumerate(roles):
    col = i % 4
    row = i // 4
    x = 1.2 + col * 3.0
    y = 2.1 + row * 0.9
    tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(x), Inches(y), Inches(2.8), Inches(0.5))
    tag.fill.solid()
    tag.fill.fore_color.rgb = WHITE
    tag.line.color.rgb = RGBColor(0xCC, 0xDD, 0xFF)
    add_textbox(slide, Inches(x), Inches(y + 0.05), Inches(2.8), Inches(0.4),
                role, 13, color=ZHIBLUE, align=PP_ALIGN.CENTER)

add_textbox(slide, Inches(0), Inches(5.8), Inches(13.333), Inches(0.4),
            '连接 DeepSeek 和知乎直答，越用越聪明', 13, color=LIGHT_GRAY)

# ========== 第5页：创作流程 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, Inches(0), Inches(0.8), Inches(13.333), Inches(0.9),
            '人人都是创作者', 42, bold=True, color=ZHIBLUE)

steps = [
    ('🎮', '玩故事', '进入一个场景\n扮演一个角色'),
    ('❤️', '爱故事', '标记精彩火花\n收藏心动瞬间'),
    ('✍️', '写短故事', '串联火花\n生成剧本'),
    ('📖', '共创长故事', '多人协作\n连载长篇'),
]
for i, (icon, title, desc) in enumerate(steps):
    x = 1.0 + i * 3.1
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     Inches(x + 0.6), Inches(2.4), Inches(1.2), Inches(1.2))
    circle.fill.solid()
    circle.fill.fore_color.rgb = WHITE
    circle.line.color.rgb = ZHIBLUE
    circle.line.width = Pt(2)
    add_textbox(slide, Inches(x), Inches(2.65), Inches(2.4), Inches(0.7),
                icon, 28, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(x), Inches(3.8), Inches(2.4), Inches(0.4),
                title, 16, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(x + 0.2), Inches(4.2), Inches(2.0), Inches(0.6),
                desc, 11, color=GRAY, align=PP_ALIGN.CENTER)
    if i < 3:
        add_textbox(slide, Inches(x + 2.3), Inches(2.8), Inches(0.6), Inches(0.5),
                    '→', 24, color=ZHIBLUE, align=PP_ALIGN.CENTER)

add_textbox(slide, Inches(0), Inches(5.5), Inches(13.333), Inches(0.4),
            '低门槛创作，高价值产出。每一个想法都值得被看见。', 13, color=LIGHT_GRAY)

# ========== 第6页：双人对白 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, Inches(0), Inches(0.8), Inches(13.333), Inches(0.9),
            '碰撞真实火花', 42, bold=True, color=ZHIBLUE)
add_textbox(slide, Inches(0), Inches(1.5), Inches(13.333), Inches(0.4),
            '两个真人 + 一个AI催化剂', 16, color=GRAY)

bubble1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(1.5), Inches(2.5), Inches(5.0), Inches(1.3))
bubble1.fill.solid()
bubble1.fill.fore_color.rgb = BG_GRAY
bubble1.line.color.rgb = RGBColor(0xE0, 0xE4, 0xEA)
add_textbox(slide, Inches(1.7), Inches(2.6), Inches(4.6), Inches(0.3),
            '医生', 11, bold=True, color=ZHIBLUE, align=PP_ALIGN.LEFT)
add_textbox(slide, Inches(1.7), Inches(2.9), Inches(4.6), Inches(0.6),
            '先推肾上腺素，准备除颤仪！', 14, color=DARK, align=PP_ALIGN.LEFT)

bubble2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(6.8), Inches(4.0), Inches(5.0), Inches(1.3))
bubble2.fill.solid()
bubble2.fill.fore_color.rgb = RGBColor(0xF0, 0xF5, 0xFF)
bubble2.line.color.rgb = RGBColor(0xCC, 0xDD, 0xFF)
add_textbox(slide, Inches(7.0), Inches(4.1), Inches(4.6), Inches(0.3),
            '家属', 11, bold=True, color=ZHIBLUE, align=PP_ALIGN.LEFT)
add_textbox(slide, Inches(7.0), Inches(4.4), Inches(4.6), Inches(0.6),
            '他是不是没有希望了？', 14, color=DARK, align=PP_ALIGN.LEFT)

bubble3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(1.5), Inches(5.5), Inches(5.0), Inches(1.3))
bubble3.fill.solid()
bubble3.fill.fore_color.rgb = BG_GRAY
bubble3.line.color.rgb = RGBColor(0xE0, 0xE4, 0xEA)
add_textbox(slide, Inches(1.7), Inches(5.6), Inches(4.6), Inches(0.3),
            '刘看山 · 催化剂', 11, bold=True, color=ZHIBLUE, align=PP_ALIGN.LEFT)
add_textbox(slide, Inches(1.7), Inches(5.9), Inches(4.6), Inches(0.6),
            '家属此刻最害怕的是什么？', 14, color=DARK, align=PP_ALIGN.LEFT)

add_textbox(slide, Inches(0), Inches(7.0), Inches(13.333), Inches(0.3),
            '这不是AI编的故事，这是真实的人说出来的话。', 12, color=LIGHT_GRAY)

# ========== 第7页：结束页 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, ZHIBLUE)
add_textbox(slide, Inches(0), Inches(2.0), Inches(13.333), Inches(1.0),
            '群像·星火', 54, bold=True, color=WHITE)
add_textbox(slide, Inches(0), Inches(3.1), Inches(13.333), Inches(0.5),
            '让真实发光，让思想变现', 20, color=RGBColor(0xDD, 0xDD, 0xDD))
url_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(4.8), Inches(4.0), Inches(3.7), Inches(0.5))
url_box.fill.solid()
url_box.fill.fore_color.rgb = RGBColor(0x00, 0x55, 0xDD)
url_box.line.color.rgb = RGBColor(0x33, 0x77, 0xEE)
add_textbox(slide, Inches(4.8), Inches(4.05), Inches(3.7), Inches(0.4),
            'http://81.70.59.228', 13, color=RGBColor(0xCC, 0xCC, 0xCC), align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(0), Inches(5.0), Inches(13.333), Inches(0.4),
            '希望在故事里，遇见你。', 14, color=RGBColor(0xAA, 0xAA, 0xAA))

# 保存
output_path = os.path.join(os.path.dirname(__file__), '..', 'docs', 'qunxiang.pptx')
os.makedirs(os.path.dirname(output_path), exist_ok=True)
prs.save(output_path)
print("PPT saved to: " + output_path)
