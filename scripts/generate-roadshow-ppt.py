#!/usr/bin/env python3
"""
群像·星火 路演 PPT 生成器
使用 python-pptx 生成 .pptx 文件
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 品牌色
COLOR_PRIMARY = RGBColor(0xFF, 0x45, 0x00)      # 星火橙
COLOR_DARK = RGBColor(0x1A, 0x1A, 0x2E)         # 深夜蓝
COLOR_LIGHT = RGBColor(0xFF, 0xFF, 0xFF)        # 白色
COLOR_GRAY = RGBColor(0x66, 0x66, 0x66)         # 灰色
COLOR_ACCENT = RGBColor(0xFF, 0xD7, 0x00)       # 金色

def add_cover_slide(prs):
    """封面页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_DARK

    # 大标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    tf = title_box.text_frame
    tf.text = "群像·星火"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(72)
    p.font.color.rgb = COLOR_PRIMARY
    p.font.bold = True
    p.font.name = "Microsoft YaHei"

    # 副标题
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(1))
    tf = sub_box.text_frame
    tf.text = "基于真实职业经验的多人协同创作平台"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(28)
    p.font.color.rgb = COLOR_LIGHT
    p.font.name = "Microsoft YaHei"

    # Slogan
    slogan_box = slide.shapes.add_textbox(Inches(1), Inches(5.3), Inches(8), Inches(0.8))
    tf = slogan_box.text_frame
    tf.text = "让真实的人，在真实的情境中，碰撞出真实的火花"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(18)
    p.font.color.rgb = COLOR_ACCENT
    p.font.name = "Microsoft YaHei"

    # 装饰线
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3), Inches(6.2), Inches(4), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_PRIMARY
    line.line.fill.background()

    return slide

def add_slide(prs, title_text, content_text):
    """添加一页标准幻灯片"""
    slide_layout = prs.slide_layouts[1]  # 标题和内容
    slide = prs.slides.add_slide(slide_layout)

    # 设置背景色
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_DARK

    # 标题
    if slide.shapes.title:
        title = slide.shapes.title
        title.text = title_text
        for paragraph in title.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.LEFT
            for run in paragraph.runs:
                run.font.name = "Microsoft YaHei"
                run.font.size = Pt(36)
                run.font.color.rgb = COLOR_PRIMARY
                run.font.bold = True

    # 内容
    if len(slide.placeholders) > 1:
        body = slide.placeholders[1]
        body.text = content_text
        tf = body.text_frame
        tf.word_wrap = True
        for paragraph in tf.paragraphs:
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.space_after = Pt(14)
            for run in paragraph.runs:
                run.font.name = "Microsoft YaHei"
                run.font.size = Pt(20)
                run.font.color.rgb = COLOR_LIGHT

    return slide

def add_end_slide(prs):
    """结尾页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_DARK

    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(1.5))
    tf = title_box.text_frame
    tf.text = "最好的故事"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(56)
    p.font.color.rgb = COLOR_PRIMARY
    p.font.bold = True
    p.font.name = "Microsoft YaHei"

    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(1))
    tf = sub_box.text_frame
    tf.text = "不是一个人关在房间里写出来的"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(28)
    p.font.color.rgb = COLOR_LIGHT
    p.font.name = "Microsoft YaHei"

    slogan_box = slide.shapes.add_textbox(Inches(1), Inches(5.2), Inches(8), Inches(0.8))
    tf = slogan_box.text_frame
    tf.text = "而是让真实的人在真实的情境中碰撞出来的"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(22)
    p.font.color.rgb = COLOR_ACCENT
    p.font.name = "Microsoft YaHei"

    thanks_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.8))
    tf = thanks_box.text_frame
    tf.text = "谢谢大家"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(32)
    p.font.color.rgb = COLOR_PRIMARY
    p.font.bold = True
    p.font.name = "Microsoft YaHei"

    return slide

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 1. 封面
    add_cover_slide(prs)

    # 2. 开场痛点
    add_slide(prs,
        "开场：一个创作者的真实困境",
        "• 创作者写剧本时，经常卡在专业细节的真实性上\n"
        "  ——急诊抢救流程、律师质证技巧、外卖员跑单逻辑\n\n"
        "• 单人创作视角永远单一\n"
        "  ——编剧采访了医生，但没采访护士、患者家属、保洁阿姨\n\n"
        "• 同一情境在不同职业眼中，是完全不同的故事\n\n"
        "• 有真实职业经验的普通人，有故事但缺乏表达渠道")

    # 3. 产品定位
    add_slide(prs,
        "群像·星火是什么？",
        "基于真实职业经验的多人协同创作平台\n\n"
        "核心理念：\n"
        "让不同职业背景的普通人，被同时扔进同一个冲突情境\n"
        "用各自的职业本能碰撞出火花\n"
        "共同完成一部一个人永远写不出的故事\n\n"
        "不是 AI 替代人类创作\n"
        "而是让真实的人、真实的经验、真实的碰撞成为创作原材料")

    # 4. 单人模式
    add_slide(prs,
        "模式一：单人灵感积累",
        "选择身份 → 浏览脑洞卡片（左滑跳过 / 右滑收藏）\n\n"
        "AI 催化引擎根据内容和身份生成引导问题：\n"
        "  「作为医生，你首先会关注哪些生命体征？」\n\n"
        "语音或文字给出反应，自动存入个人素材库\n\n"
        "适用场景：日常灵感积累、碎片化创作")

    # 5. 双人模式
    add_slide(prs,
        "模式二：双人即兴碰撞",
        "两个选择不同身份的用户，右滑收藏同一脑洞 → 进入匹配池\n\n"
        "60 秒匹配 → 实时对白室\n"
        "  一边是「急诊科医生」，一边是「患者家属」\n\n"
        "即时对话，随时标记「火花」（精彩对白片段）\n\n"
        "对话结束后：火花墙回顾 + AI 串联成完整剧本对白\n\n"
        "适用场景：即兴碰撞、快速产出对白片段")

    # 6. 多人模式
    add_slide(prs,
        "模式三：多人剧本共创",
        "导演创建副本 → 参与者认领角色\n\n"
        "导演控场推进剧情：\n"
        "  • 暂停让大家思考\n"
        "  • 发起投票决定剧情走向\n"
        "  • 喊「杀青」结束创作\n\n"
        "投票选中的精彩选项自动归档到灵感库\n"
        "AI 最终串联成一部群像故事\n\n"
        "适用场景：完整剧本创作、团队协作")

    # 7. 技术架构
    add_slide(prs,
        "技术架构：为什么我们能做出来",
        "实时协作：Socket.io 自定义 server.ts\n"
        "  Next.js + WebSocket 毫秒级同步\n\n"
        "AI 双引擎：\n"
        "  • 催化引擎：DeepSeek API 生成引导问题\n"
        "  • 串联引擎：将「火花」串联成完整故事\n\n"
        "三级降级策略：\n"
        "  DeepSeek API → 本地题库 → 通用提示\n"
        "  确保 Demo 在任何网络环境下可用\n\n"
        "TDD 测试覆盖：216 个测试全部通过")

    # 8. 商业模式
    add_slide(prs,
        "商业模式",
        "第一层：C 端免费增值\n"
        "  基础玩法免费，高级功能付费\n"
        "  AI 长格式输出 / 导演高级控场工具包\n\n"
        "第二层：B 端内容采购\n"
        "  微短剧公司、互动小说平台直接采购优质对白\n"
        "  按字数或按片段付费\n\n"
        "第三层：IP 共创分润\n"
        "  优质群像故事共同孵化 IP\n"
        "  改编为剧本、短剧、有声书，按贡献度分润")

    # 9. 里程碑
    add_slide(prs,
        "里程碑与成果",
        "Phase 1：匹配引擎\n"
        "  内存匹配池 + 优先级算法 + 60秒超时检测\n\n"
        "Phase 2：房间管理 API\n"
        "  消息 / 火花 / 暂停 / 恢复 / 结束\n\n"
        "Phase 3：实时通信 + AI 故事串联\n"
        "  WebSocket + DeepSeek API 接入\n\n"
        "Phase 4：TDD 全覆盖 + AI 催化\n"
        "  216 个测试全部通过，23 个测试文件\n"
        "  AI 催化提示生成器完成\n\n"
        "下一步：知乎 API 接入（5/9-12）")

    # 10. 结尾
    add_end_slide(prs)

    output_path = "docs/群像星火-路演PPT.pptx"
    prs.save(output_path)
    print(f"[OK] PPT generated: {output_path}")
    print(f"     Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
