from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Color constants
PRIMARY = RGBColor(0, 102, 255)
PRIMARY_LIGHT = RGBColor(61, 139, 255)
ACCENT = RGBColor(0, 212, 170)
BG_DARK = RGBColor(10, 22, 40)
BG_DARK_2 = RGBColor(13, 31, 60)
BG_LIGHT = RGBColor(255, 255, 255)
BG_GRAY = RGBColor(240, 244, 248)
TEXT_PRIMARY = RGBColor(26, 26, 46)
TEXT_SECONDARY = RGBColor(100, 116, 139)
TEXT_TERTIARY = RGBColor(148, 163, 184)
WHITE = RGBColor(255, 255, 255)
DARK_TEXT_50 = RGBColor(128, 128, 128)
DARK_TEXT_25 = RGBColor(192, 192, 192)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_gradient_bg(slide, c1, c2):
    bg = slide.background
    fill = bg.fill
    fill.gradient()
    fill.gradient_angle = 160
    fill.gradient_stops[0].color.rgb = c1
    fill.gradient_stops[1].color.rgb = c2

def add_textbox(slide, left, top, width, height, text, font_size, font_color, bold=False, align=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_rounded_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_line(slide, left, top, width, color, thickness=Pt(2)):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, thickness)
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line

def rgba(r, g, b, a):
    """Create RGBColor with approximate alpha (pptx doesn't support true alpha)"""
    # Blend with white background for light slides, black for dark
    return RGBColor(r, g, b)

W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

# ===== Slide 1: Cover =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide, BG_DARK, BG_DARK_2)
add_textbox(slide, Inches(9.5), Inches(0.5), Inches(3), Inches(0.4), '知乎黑客松 2026', Pt(11), DARK_TEXT_50, align=PP_ALIGN.RIGHT)
add_textbox(slide, Inches(0), Inches(2.2), W, Inches(0.3), 'QUNXIANG XINGHUO', Pt(12), PRIMARY_LIGHT, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(0), Inches(2.6), W, Inches(1.2), '群像·星火', Pt(72), WHITE, bold=True, align=PP_ALIGN.CENTER)
add_line(slide, Inches(5.5), Inches(3.8), Inches(2.3), PRIMARY, Pt(3))
add_textbox(slide, Inches(0), Inches(4.1), W, Inches(0.5), '让真实发光，让思想变现', Pt(24), DARK_TEXT_50, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(0), Inches(4.7), W, Inches(0.3), '一个关于创作、连接与被看见的故事', Pt(13), DARK_TEXT_25, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(0.5), Inches(6.8), Inches(6), Inches(0.3), 'TDD v4.0  |  216 Tests Passed  |  Next.js 16 + DeepSeek AI', Pt(10), DARK_TEXT_25)
add_textbox(slide, Inches(9), Inches(6.8), Inches(3.8), Inches(0.3), 'github.com/qunxiang-xinghuo', Pt(10), DARK_TEXT_25, align=PP_ALIGN.RIGHT)
add_textbox(slide, Inches(12.5), Inches(6.8), Inches(0.5), Inches(0.3), '1/10', Pt(10), DARK_TEXT_25, align=PP_ALIGN.RIGHT)

# ===== Slide 2: Pain Points =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_LIGHT)
add_textbox(slide, Inches(0), Inches(0.7), W, Inches(0.3), 'THE PROBLEM', Pt(11), PRIMARY, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(0), Inches(1.1), W, Inches(0.7), '创作者的真实困境', Pt(44), TEXT_PRIMARY, bold=True, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(3), Inches(1.8), Inches(7.3), Inches(0.4), '三个痛点，每一个都在扼杀创作的可能性', Pt(16), TEXT_SECONDARY, align=PP_ALIGN.CENTER)

pain_data = [
    ('🎭', '专业细节失真', '医生角色开口就是"先推肾上腺素"，护士只会点头——缺乏真实职业经验的支撑', '73%', '% 创作者卡在专业细节'),
    ('👁️', '视角永远单一', '编剧采访了医生，却没采访护士、家属、保洁——同一场抢救，完全是不同的故事', '1', '个视角 vs 真实世界的 N 个'),
    ('🔒', '创作门槛太高', '普通人有好故事却不知道怎么写，专业工具复杂，协作成本高昂', '90%', '% 好故事从未被写出来'),
]
for i, (icon, title, desc, stat, stat_label) in enumerate(pain_data):
    left = Inches(0.8 + i * 4.2)
    card = add_rounded_rect(slide, left, Inches(2.6), Inches(3.8), Inches(4.2), BG_LIGHT, RGBColor(0, 102, 255))
    card.line.color.rgb = RGBColor(220, 230, 245)
    add_textbox(slide, left, Inches(2.9), Inches(3.8), Inches(0.6), icon, Pt(36), TEXT_PRIMARY, align=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.2), Inches(3.5), Inches(3.4), Inches(0.4), title, Pt(18), TEXT_PRIMARY, bold=True, align=PP_ALIGN.CENTER)
    tb = add_textbox(slide, left + Inches(0.2), Inches(4.0), Inches(3.4), Inches(1.2), desc, Pt(12), TEXT_SECONDARY, align=PP_ALIGN.CENTER)
    tb.text_frame.paragraphs[0].line_spacing = 1.4
    add_textbox(slide, left, Inches(5.3), Inches(3.8), Inches(0.5), stat, Pt(32), PRIMARY, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, left, Inches(5.7), Inches(3.8), Inches(0.3), stat_label, Pt(10), TEXT_TERTIARY, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(12.5), Inches(6.8), Inches(0.5), Inches(0.3), '2/10', Pt(10), TEXT_TERTIARY, align=PP_ALIGN.RIGHT)

# ===== Slide 3: Solution =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide, BG_DARK, BG_DARK_2)
add_textbox(slide, Inches(0), Inches(0.7), W, Inches(0.3), 'OUR SOLUTION', Pt(11), PRIMARY_LIGHT, align=PP_ALIGN.CENTER)
quote = '让不同职业背景的普通人，被同时扔进同一个冲突情境，用各自的职业本能碰撞出火花'
add_textbox(slide, Inches(1.5), Inches(2.0), Inches(10.3), Inches(1.0), quote, Pt(28), WHITE, bold=True, align=PP_ALIGN.CENTER)
desc = '群像·星火是一个基于真实职业经验的多人协同创作平台。\n不是 AI 替代人类创作，而是让真实的人、真实的经验、真实的碰撞成为创作的原材料。'
tb = add_textbox(slide, Inches(2.5), Inches(3.4), Inches(8.3), Inches(1.0), desc, Pt(15), DARK_TEXT_50, align=PP_ALIGN.CENTER)
tb.text_frame.paragraphs[0].line_spacing = 1.6
add_textbox(slide, Inches(12.5), Inches(6.8), Inches(0.5), Inches(0.3), '3/10', Pt(10), DARK_TEXT_25, align=PP_ALIGN.RIGHT)

# ===== Slide 4: Three Modes =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_GRAY)
add_textbox(slide, Inches(0), Inches(0.7), W, Inches(0.3), 'PRODUCT', Pt(11), PRIMARY, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(0), Inches(1.1), W, Inches(0.7), '三个模式，一个信念', Pt(44), TEXT_PRIMARY, bold=True, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(3), Inches(1.8), Inches(7.3), Inches(0.4), '覆盖从日常灵感积累到完整剧本创作的全流程', Pt(16), TEXT_SECONDARY, align=PP_ALIGN.CENTER)

modes = [
    ('MODE 01', '🤖', '人机模式', '和刘看山聊聊你的脑洞\nAI 催化引擎生成引导问题\n一个永远在线的倾听者'),
    ('MODE 02', '💬', '双人模式', '两个真人 + 一个 AI 催化剂\n围绕同一个情境即时对话\n60 秒匹配，碰撞真实火花'),
    ('MODE 03', '🎭', '故事模式', '选角色，揭开历史谜题\n导演控场，四幕推进\n真相由你揭晓'),
]
for i, (num, icon, title, desc) in enumerate(modes):
    left = Inches(0.8 + i * 4.2)
    card = add_rounded_rect(slide, left, Inches(2.5), Inches(3.8), Inches(4.3), BG_LIGHT, RGBColor(220, 230, 245))
    add_line(slide, left, Inches(2.5), Inches(3.8), PRIMARY, Pt(3))
    add_textbox(slide, left, Inches(2.7), Inches(3.8), Inches(0.3), num, Pt(11), PRIMARY, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, left, Inches(3.1), Inches(3.8), Inches(0.5), icon, Pt(32), TEXT_PRIMARY, align=PP_ALIGN.CENTER)
    add_textbox(slide, left, Inches(3.6), Inches(3.8), Inches(0.4), title, Pt(20), TEXT_PRIMARY, bold=True, align=PP_ALIGN.CENTER)
    tb = add_textbox(slide, left + Inches(0.2), Inches(4.1), Inches(3.4), Inches(1.5), desc, Pt(12), TEXT_SECONDARY, align=PP_ALIGN.CENTER)
    tb.text_frame.paragraphs[0].line_spacing = 1.5
add_textbox(slide, Inches(12.5), Inches(6.8), Inches(0.5), Inches(0.3), '4/10', Pt(10), TEXT_TERTIARY, align=PP_ALIGN.RIGHT)

# ===== Slide 5: Demo =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide, BG_DARK, BG_DARK_2)
add_textbox(slide, Inches(0), Inches(0.7), W, Inches(0.3), 'LIVE DEMO', Pt(11), PRIMARY_LIGHT, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(0), Inches(1.1), W, Inches(0.6), '碰撞真实火花', Pt(40), WHITE, bold=True, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(0), Inches(1.7), W, Inches(0.3), '两个真人 + 一个 AI 催化剂', Pt(15), DARK_TEXT_50, align=PP_ALIGN.CENTER)

mock_left = Inches(3.8)
mock_top = Inches(2.3)
mock_w = Inches(5.7)
mock_h = Inches(3.8)
frame = add_rounded_rect(slide, mock_left, mock_top, mock_w, mock_h, RGBColor(25, 35, 55), RGBColor(50, 60, 80))

dot_colors = [RGBColor(255, 95, 87), RGBColor(254, 188, 46), RGBColor(40, 200, 64)]
for i, c in enumerate(dot_colors):
    d = slide.shapes.add_shape(MSO_SHAPE.OVAL, mock_left + Inches(0.3 + i * 0.25), mock_top + Inches(0.2), Inches(0.12), Inches(0.12))
    d.fill.solid(); d.fill.fore_color.rgb = c; d.line.fill.background()
add_textbox(slide, mock_left + Inches(1.0), mock_top + Inches(0.15), Inches(3), Inches(0.2), '深夜急诊室 · 双人对话', Pt(9), DARK_TEXT_25)

bubbles = [
    (mock_left + Inches(0.3), mock_top + Inches(0.6), Inches(4.5), Inches(0.7), '医生', '先推肾上腺素，准备除颤仪！患者室颤了！', RGBColor(40, 50, 70), RGBColor(220, 225, 235)),
    (mock_left + Inches(1.0), mock_top + Inches(1.4), Inches(4.5), Inches(0.6), '家属', '他是不是……没有希望了？', RGBColor(0, 60, 140), RGBColor(220, 230, 245)),
    (mock_left + Inches(1.2), mock_top + Inches(2.1), Inches(3.3), Inches(0.7), '刘看山 · 催化剂', '家属此刻最害怕的是什么？', RGBColor(0, 80, 180), WHITE),
]
for left, top, w, h, label, text, fill_c, text_c in bubbles:
    bubble = add_rounded_rect(slide, left, top, w, h, fill_c, None)
    lc = PRIMARY_LIGHT if '催化剂' not in label else ACCENT
    add_textbox(slide, left + Inches(0.15), top + Inches(0.08), w - Inches(0.3), Inches(0.2), label, Pt(8), lc)
    add_textbox(slide, left + Inches(0.15), top + Inches(0.25), w - Inches(0.3), Inches(0.4), text, Pt(11), text_c)

add_textbox(slide, Inches(0), Inches(6.3), W, Inches(0.3), '这不是 AI 编的故事，这是真实的人说出来的话', Pt(12), DARK_TEXT_25, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(12.5), Inches(6.8), Inches(0.5), Inches(0.3), '5/10', Pt(10), DARK_TEXT_25, align=PP_ALIGN.RIGHT)

# ===== Slide 6: Agent =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_LIGHT)
add_textbox(slide, Inches(0), Inches(0.6), W, Inches(0.3), 'AI AGENT', Pt(11), PRIMARY, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(0), Inches(1.0), W, Inches(0.6), '一个 Agent，11 种角色', Pt(44), TEXT_PRIMARY, bold=True, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(0), Inches(1.6), W, Inches(0.3), '刘看山 —— 最懂你的 AI 副导演', Pt(15), TEXT_SECONDARY, align=PP_ALIGN.CENTER)

roles = [
    ('Companion', '陪伴者', PRIMARY),
    ('DM', '主持人', ACCENT),
    ('Story Fallback', '故事替补', RGBColor(245, 158, 11)),
    ('AI 副导演', '剧情把控', RGBColor(139, 92, 246)),
    ('Catalyst', '催化剂', RGBColor(239, 68, 68)),
    ('Healer', '疗愈师', RGBColor(236, 72, 153)),
    ('Reviewer', '审稿人', RGBColor(20, 184, 166)),
    ('Summarizer', '总结者', RGBColor(99, 102, 241)),
    ('Knowledge', '知识投喂', RGBColor(249, 115, 22)),
    ('Mediator', '调解员', RGBColor(6, 182, 212)),
    ('Creative', '创作助手', RGBColor(132, 204, 22)),
]
for i, (name, desc, color) in enumerate(roles):
    col = i % 4
    row = i // 4
    left = Inches(0.9 + col * 3.1)
    top = Inches(2.2 + row * 0.9)
    d = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top + Inches(0.08), Inches(0.18), Inches(0.18))
    d.fill.solid(); d.fill.fore_color.rgb = color; d.line.fill.background()
    add_textbox(slide, left + Inches(0.28), top, Inches(1.4), Inches(0.35), name, Pt(13), TEXT_PRIMARY, bold=True)
    add_textbox(slide, left + Inches(1.7), top, Inches(1.2), Inches(0.35), desc, Pt(11), TEXT_TERTIARY)
add_textbox(slide, Inches(12.5), Inches(6.8), Inches(0.5), Inches(0.3), '6/10', Pt(10), TEXT_TERTIARY, align=PP_ALIGN.RIGHT)

# ===== Slide 7: Tech =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide, BG_DARK, BG_DARK_2)
add_textbox(slide, Inches(0), Inches(0.6), W, Inches(0.3), 'TECHNOLOGY', Pt(11), PRIMARY_LIGHT, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(0), Inches(1.0), W, Inches(0.6), '四层架构，全栈自研', Pt(40), WHITE, bold=True, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(0), Inches(1.6), W, Inches(0.3), '从数据库到 AI 双引擎，每一层都经过 216 个测试验证', Pt(14), DARK_TEXT_25, align=PP_ALIGN.CENTER)

layers = [
    ('01', '前端', 'Next.js 16 + React 19 + Tailwind v4'),
    ('02', 'API', '22 个路由 + NextAuth + Zod 验证'),
    ('03', '业务', 'Socket.io + 匹配引擎 + 房间管理'),
    ('04', '数据', 'Prisma 7 + SQLite + 三级降级策略'),
]
for i, (num, title, desc) in enumerate(layers):
    top = Inches(2.2 + i * 0.95)
    layer = add_rounded_rect(slide, Inches(0.8), top, Inches(6.5), Inches(0.75), RGBColor(25, 35, 55), RGBColor(50, 60, 80))
    add_textbox(slide, Inches(1.0), top + Inches(0.15), Inches(0.6), Inches(0.4), num, Pt(11), PRIMARY_LIGHT, bold=True)
    add_textbox(slide, Inches(1.5), top + Inches(0.15), Inches(1.0), Inches(0.4), title, Pt(13), WHITE, bold=True)
    add_textbox(slide, Inches(2.6), top + Inches(0.18), Inches(4.5), Inches(0.4), desc, Pt(11), DARK_TEXT_50)

metrics = [('216', 'TESTS PASSED', PRIMARY_LIGHT), ('23', 'TEST FILES', ACCENT), ('22', 'API ROUTES', WHITE), ('0', 'FAILED', WHITE)]
for i, (val, label, color) in enumerate(metrics):
    left = Inches(8.0 + (i % 2) * 2.5)
    top = Inches(2.2 + (i // 2) * 2.0)
    card = add_rounded_rect(slide, left, top, Inches(2.2), Inches(1.6), RGBColor(25, 35, 55), RGBColor(50, 60, 80))
    add_textbox(slide, left, top + Inches(0.3), Inches(2.2), Inches(0.6), val, Pt(36), color, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, left, top + Inches(0.85), Inches(2.2), Inches(0.3), label, Pt(10), DARK_TEXT_25, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(12.5), Inches(6.8), Inches(0.5), Inches(0.3), '7/10', Pt(10), DARK_TEXT_25, align=PP_ALIGN.RIGHT)

# ===== Slide 8: Business =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_GRAY)
add_textbox(slide, Inches(0), Inches(0.7), W, Inches(0.3), 'BUSINESS MODEL', Pt(11), PRIMARY, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(0), Inches(1.1), W, Inches(0.6), '三层变现模型', Pt(44), TEXT_PRIMARY, bold=True, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(3), Inches(1.7), Inches(7.3), Inches(0.4), '从免费增值到 IP 共创，每一步都有清晰的商业化路径', Pt(15), TEXT_SECONDARY, align=PP_ALIGN.CENTER)

biz = [
    ('C 端', PRIMARY, '🎯', '免费增值', '基础玩法免费，高级功能付费\nAI 长格式输出 · 导演高级控场包'),
    ('B 端', ACCENT, '🏢', '内容采购', '微短剧 / 互动小说平台\n按字数或片段采购优质对白'),
    ('IP', RGBColor(245, 158, 11), '💎', '共创分润', '优质群像故事共同孵化 IP\n剧本 / 短剧 / 有声书 · 按贡献分润'),
]
for i, (badge, badge_color, icon, title, desc) in enumerate(biz):
    left = Inches(0.8 + i * 4.2)
    card = add_rounded_rect(slide, left, Inches(2.5), Inches(3.8), Inches(4.0), BG_LIGHT, RGBColor(220, 230, 245))
    badge_shape = add_rounded_rect(slide, left + Inches(1.3), Inches(2.2), Inches(1.2), Inches(0.35), badge_color, None)
    add_textbox(slide, left + Inches(1.3), Inches(2.22), Inches(1.2), Inches(0.3), badge, Pt(10), WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, left, Inches(2.8), Inches(3.8), Inches(0.5), icon, Pt(28), TEXT_PRIMARY, align=PP_ALIGN.CENTER)
    add_textbox(slide, left, Inches(3.3), Inches(3.8), Inches(0.4), title, Pt(18), TEXT_PRIMARY, bold=True, align=PP_ALIGN.CENTER)
    tb = add_textbox(slide, left + Inches(0.2), Inches(3.75), Inches(3.4), Inches(1.2), desc, Pt(12), TEXT_SECONDARY, align=PP_ALIGN.CENTER)
    tb.text_frame.paragraphs[0].line_spacing = 1.5
add_textbox(slide, Inches(12.5), Inches(6.8), Inches(0.5), Inches(0.3), '8/10', Pt(10), TEXT_TERTIARY, align=PP_ALIGN.RIGHT)

# ===== Slide 9: Milestones =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide, BG_DARK, BG_DARK_2)
add_textbox(slide, Inches(0), Inches(0.6), W, Inches(0.3), 'MILESTONES', Pt(11), PRIMARY_LIGHT, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(0), Inches(1.0), W, Inches(0.6), '四周，从零到完整 MVP', Pt(40), WHITE, bold=True, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(0), Inches(1.6), W, Inches(0.3), '严格的 TDD 开发流程，每一阶段都有完整的测试覆盖', Pt(14), DARK_TEXT_25, align=PP_ALIGN.CENTER)

add_line(slide, Inches(1.8), Inches(2.6), Inches(9.7), PRIMARY, Pt(2))

timeline = [
    ('PHASE 1', '匹配引擎', '9 tests', PRIMARY),
    ('PHASE 2', '房间管理', '16 tests', PRIMARY),
    ('PHASE 3', 'WebSocket + AI', '21 tests', PRIMARY),
    ('PHASE 4', 'TDD 全覆盖', '216 tests', ACCENT),
]
for i, (phase, title, tests, color) in enumerate(timeline):
    left = Inches(1.0 + i * 2.8)
    d = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.8), Inches(2.45), Inches(0.3), Inches(0.3))
    d.fill.solid(); d.fill.fore_color.rgb = BG_DARK; d.line.color.rgb = color; d.line.width = Pt(3)
    add_textbox(slide, left, Inches(2.9), Inches(2.4), Inches(0.25), phase, Pt(10), PRIMARY_LIGHT, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, left, Inches(3.2), Inches(2.4), Inches(0.3), title, Pt(14), WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, left, Inches(3.5), Inches(2.4), Inches(0.25), tests, Pt(11), DARK_TEXT_25, align=PP_ALIGN.CENTER)

bignums = [('216', 'Tests', PRIMARY_LIGHT), ('22', 'APIs', ACCENT), ('11', 'Personas', WHITE), ('3', 'Modes', WHITE)]
for i, (num, label, color) in enumerate(bignums):
    left = Inches(1.5 + i * 3.0)
    add_textbox(slide, left, Inches(4.3), Inches(2.5), Inches(0.8), num, Pt(52), color, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, left, Inches(5.0), Inches(2.5), Inches(0.25), label, Pt(11), DARK_TEXT_25, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(12.5), Inches(6.8), Inches(0.5), Inches(0.3), '9/10', Pt(10), DARK_TEXT_25, align=PP_ALIGN.RIGHT)

# ===== Slide 10: End =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide, BG_DARK, BG_DARK_2)
add_textbox(slide, Inches(0), Inches(2.0), W, Inches(1.0), '群像·星火', Pt(72), WHITE, bold=True, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(0), Inches(3.0), W, Inches(0.4), '让真实发光，让思想变现', Pt(22), DARK_TEXT_50, align=PP_ALIGN.CENTER)
url_box = add_rounded_rect(slide, Inches(4.8), Inches(3.7), Inches(3.7), Inches(0.5), RGBColor(25, 35, 55), RGBColor(60, 70, 90))
add_textbox(slide, Inches(4.8), Inches(3.75), Inches(3.7), Inches(0.4), 'http://81.70.59.228', Pt(13), DARK_TEXT_50, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(0), Inches(4.5), W, Inches(0.3), '希望在故事里，遇见你。', Pt(14), DARK_TEXT_25, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(12.5), Inches(6.8), Inches(0.5), Inches(0.3), '10/10', Pt(10), DARK_TEXT_25, align=PP_ALIGN.RIGHT)

# Save
prs.save('docs/群像星火-路演PPT-v4.pptx')
print('Generated: docs/群像星火-路演PPT-v4.pptx (10 slides)')
