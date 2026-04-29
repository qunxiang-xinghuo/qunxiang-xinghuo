#!/usr/bin/env python3
"""
群像·星火 路演 PPT 生成器 v3
更精美版本：渐变背景 + 数据可视化 + 统一页脚 + 现代排版
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 品牌色
C_ORANGE  = RGBColor(0xFF, 0x45, 0x00)
C_DARK    = RGBColor(0x13, 0x13, 0x23)
C_LIGHT   = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY    = RGBColor(0x66, 0x66, 0x77)
C_GOLD    = RGBColor(0xFF, 0xD7, 0x00)
C_CARD    = RGBColor(0x22, 0x22, 0x35)
C_GREEN   = RGBColor(0x22, 0xC5, 0x5E)
C_BLUE    = RGBColor(0x3B, 0x82, 0xF6)
C_PINK    = RGBColor(0xEC, 0x48, 0x99)
C_VIOLET  = RGBColor(0x8B, 0x5C, 0xF6)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def bg(slide, color=C_DARK):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def txt(slide, l, t, w, h, text, size=24, color=C_LIGHT, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.alignment = align
    p.word_wrap = True
    for r in p.runs:
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.name = "Microsoft YaHei"
    return box

def rect(slide, l, t, w, h, fill=None, line=None, radius=0):
    if radius:
        s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    else:
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
    else:
        s.line.fill.background()
    return s

def oval(slide, l, t, w, h, fill=None):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s

def footer(slide, page_num, total=10):
    txt(slide, 12, 7.1, 1, 0.3, f"{page_num}/{total}", 12, C_GRAY, False, PP_ALIGN.RIGHT)
    txt(slide, 0.5, 7.1, 3, 0.3, "群像·星火", 12, C_GRAY)
    rect(slide, 0.5, 7.0, 12.3, 0.01, C_GRAY)

# ===================== 第1页：封面 =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
oval(s, -2, -1, 5, 5, RGBColor(0xFF,0x45,0x00))
oval(s, 10, 5, 4, 4, RGBColor(0xFF,0x45,0x00))
oval(s, 9, 1, 1.5, 1.5, RGBColor(0xFF,0xD7,0x00))
txt(s, 1, 2.2, 11.3, 1.5, "群像·星火", 96, C_ORANGE, True, PP_ALIGN.CENTER)
txt(s, 1, 3.6, 11.3, 0.8, "基于真实职业经验的多人协同创作平台", 32, C_LIGHT, False, PP_ALIGN.CENTER)
txt(s, 1, 4.4, 11.3, 0.6, "让真实的人，在真实的情境中，碰撞出真实的火花", 20, C_GOLD, False, PP_ALIGN.CENTER)
r = rect(s, 5.5, 5.2, 2.3, 0.04, C_ORANGE)
r.line.fill.background()
txt(s, 1, 5.8, 11.3, 0.5, "TDD v4.0  |  216 Tests Passed  |  Next.js 16 + DeepSeek AI  |  Socket.io 实时通信", 16, C_GRAY, False, PP_ALIGN.CENTER)
footer(s, 1)

# ===================== 第2页：痛点 =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
txt(s, 0.5, 0.5, 5, 0.8, "创作者的困境", 44, C_ORANGE, True)
q = oval(s, 0.5, 1.8, 4, 4, C_ORANGE)
txt(s, 0.5, 2.6, 4, 1.5, "?", 160, C_DARK, True, PP_ALIGN.CENTER)
pains = [
    ("专业细节", "写不对", "急诊流程 / 律师质证 / 外卖跑单"),
    ("视角单一", "写不深", "采访了医生，没采访护士和家属"),
    ("有经历的人", "没渠道", "退休阿姨 / 急诊护士 / 程序员"),
]
for i, (t1, t2, desc) in enumerate(pains):
    y = 1.2 + i * 1.9
    card = rect(s, 5.5, y, 7, 1.5, C_CARD, radius=1)
    card.line.color.rgb = RGBColor(0x33,0x33,0x44)
    txt(s, 5.8, y+0.15, 6, 0.6, f"{t1}，{t2}", 28, C_ORANGE, True)
    txt(s, 5.8, y+0.75, 6, 0.5, desc, 16, C_GRAY)
footer(s, 2)

# ===================== 第3页：产品定位 =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
txt(s, 0.5, 0.5, 12, 0.8, "群像·星火是什么？", 44, C_ORANGE, True, PP_ALIGN.CENTER)
oval(s, 3.5, 2.0, 3, 3, C_BLUE)
txt(s, 3.5, 3.0, 3, 0.6, "医生", 28, C_LIGHT, True, PP_ALIGN.CENTER)
sp = oval(s, 5.5, 2.5, 2.2, 2.2, C_ORANGE)
txt(s, 5.5, 3.2, 2.2, 0.5, "碰撞", 20, C_DARK, True, PP_ALIGN.CENTER)
oval(s, 7, 2.0, 3, 3, C_PINK)
txt(s, 7, 3.0, 3, 0.6, "导演", 28, C_LIGHT, True, PP_ALIGN.CENTER)
txt(s, 1, 5.5, 11.3, 0.8, "不同职业的人被扔进同一冲突情境，用各自的职业本能碰撞出火花", 24, C_LIGHT, False, PP_ALIGN.CENTER)
txt(s, 1, 6.2, 11.3, 0.5, "不是 AI 替代人类，而是让真实的人、真实的经验、真实的碰撞成为创作原材料", 18, C_GOLD, False, PP_ALIGN.CENTER)
footer(s, 3)

# ===================== 第4页：单人模式 =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
txt(s, 0.5, 0.4, 12, 0.7, "模式一：单人灵感积累", 40, C_ORANGE, True, PP_ALIGN.CENTER)
steps = [
    ("1", "选择身份", "医生 / 导演 / 外卖员"),
    ("2", "浏览脑洞", "左滑跳过 · 右滑收藏"),
    ("3", "AI 催化", "DeepSeek 生成引导问题"),
    ("4", "存入素材库", "语音 / 文字反应"),
]
for i, (num, title, sub) in enumerate(steps):
    y = 1.5 + i * 1.4
    n = oval(s, 0.8, y, 0.7, 0.7, C_ORANGE)
    txt(s, 0.8, y+0.1, 0.7, 0.5, num, 20, C_DARK, True, PP_ALIGN.CENTER)
    txt(s, 1.7, y+0.05, 3, 0.5, title, 22, C_LIGHT, True)
    txt(s, 1.7, y+0.45, 3, 0.4, sub, 14, C_GRAY)
phone = rect(s, 5.5, 1.3, 3.2, 5.2, RGBColor(0x0a,0x0a,0x18), radius=1)
phone.line.color.rgb = RGBColor(0x44,0x44,0x55)
phone.line.width = Pt(3)
txt(s, 5.5, 1.6, 3.2, 0.4, "深夜急诊室", 14, C_LIGHT, True, PP_ALIGN.CENTER)
rect(s, 5.7, 2.2, 2.8, 1.4, C_CARD, radius=1)
txt(s, 5.7, 2.5, 2.8, 0.8, "外卖员因\n过度劳累晕倒", 12, C_GRAY, False, PP_ALIGN.CENTER)
rect(s, 5.7, 3.9, 2.8, 1.0, C_ORANGE, radius=1)
txt(s, 5.7, 4.15, 2.8, 0.5, "你首先会关注\n哪些生命体征？", 12, C_DARK, True, PP_ALIGN.CENTER)
oval(s, 6.8, 5.5, 0.6, 0.6, RGBColor(0x33,0x33,0x44))
txt(s, 6.8, 5.6, 0.6, 0.4, "mic", 10, C_ORANGE, False, PP_ALIGN.CENTER)
txt(s, 9.5, 1.8, 3, 1.2, "5+", 72, C_ORANGE, True, PP_ALIGN.RIGHT)
txt(s, 9.5, 2.9, 3, 0.4, "职业标签", 16, C_GRAY, False, PP_ALIGN.RIGHT)
txt(s, 9.5, 4.0, 3, 1.2, "60+", 72, C_ORANGE, True, PP_ALIGN.RIGHT)
txt(s, 9.5, 5.1, 3, 0.4, "引导问题", 16, C_GRAY, False, PP_ALIGN.RIGHT)
footer(s, 4)

# ===================== 第5页：双人模式 =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
txt(s, 0.5, 0.4, 12, 0.7, "模式二：双人即兴碰撞", 40, C_ORANGE, True, PP_ALIGN.CENTER)
oval(s, 2, 1.8, 3, 3, C_BLUE)
txt(s, 2, 2.9, 3, 0.6, "医生", 28, C_LIGHT, True, PP_ALIGN.CENTER)
txt(s, 1.5, 5.0, 4, 0.8, "先推肾上腺素！\n患者室颤了！", 16, C_LIGHT, False, PP_ALIGN.CENTER)
oval(s, 8.3, 1.8, 3, 3, C_PINK)
txt(s, 8.3, 2.9, 3, 0.6, "导演", 28, C_LIGHT, True, PP_ALIGN.CENTER)
txt(s, 7.8, 5.0, 4, 0.8, "等等，情绪高潮\n应该在后面", 16, C_LIGHT, False, PP_ALIGN.CENTER)
sp = oval(s, 5.5, 2.3, 2.2, 2.2, C_ORANGE)
txt(s, 5.5, 3.0, 2.2, 0.5, "碰撞", 20, C_DARK, True, PP_ALIGN.CENTER)
data = [("60s", "匹配等待"), ("<100ms", "WebSocket同步"), ("AI", "自动串联剧本")]
for i, (num, label) in enumerate(data):
    x = 2.5 + i * 3.5
    txt(s, x, 6.0, 3, 1, num, 40, C_ORANGE, True, PP_ALIGN.CENTER)
    txt(s, x, 6.9, 3, 0.4, label, 14, C_GRAY, False, PP_ALIGN.CENTER)
footer(s, 5)

# ===================== 第6页：多人模式 =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
txt(s, 0.5, 0.4, 12, 0.7, "模式三：多人剧本共创", 40, C_ORANGE, True, PP_ALIGN.CENTER)
director = oval(s, 5.3, 1.2, 2.8, 2.8, C_ORANGE)
txt(s, 5.3, 2.2, 2.8, 0.6, "导演", 24, C_DARK, True, PP_ALIGN.CENTER)
roles = [
    ("医生", 2.5, 1.5, C_BLUE),
    ("护士", 9.0, 1.5, C_GREEN),
    ("家属", 2.0, 3.5, C_VIOLET),
    ("律师", 9.5, 3.5, C_PINK),
]
for label, x, y, color in roles:
    p = oval(s, x, y, 1.6, 1.6, C_CARD)
    p.line.color.rgb = RGBColor(0x44,0x44,0x55)
    txt(s, x, y+0.5, 1.6, 0.4, label, 14, C_LIGHT, False, PP_ALIGN.CENTER)
controls = [("暂停", 3.5), ("投票", 5.5), ("杀青", 7.5)]
for label, x in controls:
    btn = rect(s, x, 5.0, 1.5, 0.6, C_CARD, radius=1)
    btn.line.color.rgb = C_ORANGE
    txt(s, x, 5.1, 1.5, 0.4, label, 16, C_ORANGE, False, PP_ALIGN.CENTER)
flow = ["创建副本", "认领角色", "导演控场", "投票决策", "AI串联", "署名墙"]
for i, step in enumerate(flow):
    x = 0.6 + i * 2.1
    box = rect(s, x, 6.0, 1.8, 0.45, C_CARD, radius=1)
    box.line.color.rgb = RGBColor(0x44,0x44,0x55)
    txt(s, x, 6.08, 1.8, 0.35, step, 12, C_LIGHT, False, PP_ALIGN.CENTER)
footer(s, 6)

# ===================== 第7页：技术架构 =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
txt(s, 0.5, 0.4, 12, 0.7, "技术架构", 40, C_ORANGE, True, PP_ALIGN.CENTER)
layers = [
    ("Next.js 16 + React 19 + Tailwind v4", C_BLUE, 0),
    ("API Routes + NextAuth + Zod", C_VIOLET, 0.1),
    ("Socket.io + match-engine + room-manager", C_PINK, 0.2),
    ("Prisma 7 + SQLite", C_GREEN, 0.3),
]
for i, (label, color, delay) in enumerate(layers):
    y = 1.4 + i * 1.1
    rect(s, 2.5, y, 8, 0.9, color)
    txt(s, 2.5, y+0.2, 8, 0.5, label, 20, C_LIGHT, True, PP_ALIGN.CENTER)
labels = ["前端", "API层", "业务层", "数据层"]
for i, label in enumerate(labels):
    y = 1.4 + i * 1.1
    txt(s, 0.5, y+0.2, 1.8, 0.5, label, 16, C_GRAY, False, PP_ALIGN.RIGHT)
ds = rect(s, 11, 2, 1.8, 1.8, RGBColor(0x11,0x44,0x77), radius=1)
ds.line.fill.background()
txt(s, 11, 2.6, 1.8, 0.5, "DeepSeek\nAI", 14, C_LIGHT, True, PP_ALIGN.CENTER)
big_data = [
    ("216", "Tests Passed", C_GREEN),
    ("23", "Test Files", C_ORANGE),
    ("0", "Failed", C_GREEN),
]
for i, (num, label, color) in enumerate(big_data):
    x = 1.5 + i * 4
    txt(s, x, 5.8, 3.5, 1.2, num, 80, color, True, PP_ALIGN.CENTER)
    txt(s, x, 6.8, 3.5, 0.4, label, 16, C_GRAY, False, PP_ALIGN.CENTER)
footer(s, 7)

# ===================== 第8页：商业模式 =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
txt(s, 0.5, 0.4, 12, 0.7, "商业模式", 40, C_ORANGE, True, PP_ALIGN.CENTER)
b1 = rect(s, 2, 5.0, 9.3, 1.4, C_GREEN, radius=1)
b1.line.fill.background()
txt(s, 2, 5.2, 9.3, 0.8, "C 端免费增值", 32, C_DARK, True, PP_ALIGN.CENTER)
txt(s, 2, 5.85, 9.3, 0.4, "基础免费 + AI高级功能付费", 16, C_DARK, False, PP_ALIGN.CENTER)
b2 = rect(s, 3.5, 3.2, 6.3, 1.4, C_BLUE, radius=1)
b2.line.fill.background()
txt(s, 3.5, 3.4, 6.3, 0.8, "B 端内容采购", 32, C_LIGHT, True, PP_ALIGN.CENTER)
txt(s, 3.5, 4.05, 6.3, 0.4, "微短剧公司 / 互动小说平台", 16, C_LIGHT, False, PP_ALIGN.CENTER)
b3 = rect(s, 5, 1.4, 3.3, 1.4, C_ORANGE, radius=1)
b3.line.fill.background()
txt(s, 5, 1.6, 3.3, 0.8, "IP 共创分润", 32, C_DARK, True, PP_ALIGN.CENTER)
txt(s, 5, 2.25, 3.3, 0.4, "剧本 / 短剧 / 有声书", 16, C_DARK, False, PP_ALIGN.CENTER)
footer(s, 8)

# ===================== 第9页：里程碑 =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
txt(s, 0.5, 0.4, 12, 0.7, "里程碑与成果", 40, C_ORANGE, True, PP_ALIGN.CENTER)
rect(s, 0.5, 3.3, 12.3, 0.06, C_GRAY)
phases = [
    ("Phase 1", "匹配引擎", "9 tests", 0),
    ("Phase 2", "房间管理", "16 tests", 0.1),
    ("Phase 3", "WebSocket\n+ AI串联", "21 tests", 0.2),
    ("Phase 4", "TDD全覆盖\n+ AI催化", "216 tests", 0.3),
]
for i, (label, title, tests, delay) in enumerate(phases):
    x = 1.2 + i * 3
    node = oval(s, x, 3.1, 0.4, 0.4, C_ORANGE)
    node.line.fill.background()
    txt(s, x-0.8, 2.3, 2, 0.5, label, 18, C_ORANGE, True, PP_ALIGN.CENTER)
    txt(s, x-0.8, 3.6, 2, 0.8, title, 16, C_LIGHT, False, PP_ALIGN.CENTER)
    txt(s, x-0.8, 4.4, 2, 0.4, tests, 14, C_GREEN, False, PP_ALIGN.CENTER)
txt(s, 0.5, 5.5, 12.3, 1.2, "216 / 216", 100, C_GREEN, True, PP_ALIGN.CENTER)
txt(s, 0.5, 6.5, 12.3, 0.4, "All Tests Passed  |  23 Test Files  |  22 API Routes Covered", 16, C_GRAY, False, PP_ALIGN.CENTER)
footer(s, 9)

# ===================== 第10页：结尾 =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
oval(s, -2, -2, 5, 5, RGBColor(0xFF,0x45,0x00))
oval(s, 10, 4, 4, 4, RGBColor(0xFF,0x45,0x00))
txt(s, 0.5, 2.0, 12.3, 1.2, "最好的故事", 80, C_ORANGE, True, PP_ALIGN.CENTER)
txt(s, 0.5, 3.2, 12.3, 0.7, "不是一个人关在房间里写出来的", 28, C_LIGHT, False, PP_ALIGN.CENTER)
txt(s, 0.5, 3.9, 12.3, 0.6, "而是让真实的人在真实的情境中碰撞出来的", 24, C_GOLD, False, PP_ALIGN.CENTER)
txt(s, 0.5, 5.0, 12.3, 0.8, "谢谢大家", 48, C_ORANGE, True, PP_ALIGN.CENTER)
txt(s, 0.5, 5.8, 12.3, 0.4, "github.com/qunxiang-xinghuo", 16, C_GRAY, False, PP_ALIGN.CENTER)
footer(s, 10)

output = "docs/群像星火-路演PPT-v3.pptx"
prs.save(output)
print(f"[OK] PPT v3 generated: {output}")
print(f"     Total slides: {len(prs.slides)}")
print(f"     Slide size: 16:9 (13.333 x 7.5 inches)")
