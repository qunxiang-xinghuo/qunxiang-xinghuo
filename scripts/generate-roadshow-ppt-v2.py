#!/usr/bin/env python3
"""
群像·星火 路演 PPT 生成器 v2
大图少字风格：几何图形模拟 UI + 大数字 + 流程图 + 极简文案
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 品牌色
C_ORANGE = RGBColor(0xFF, 0x45, 0x00)
C_DARK   = RGBColor(0x1A, 0x1A, 0x2E)
C_LIGHT  = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY   = RGBColor(0x88, 0x88, 0x99)
C_GOLD   = RGBColor(0xFF, 0xD7, 0x00)
C_CARD   = RGBColor(0x25, 0x25, 0x3A)
C_GREEN  = RGBColor(0x22, 0xC5, 0x5E)

def add_bg(slide, color=C_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, font_size=24, color=C_LIGHT, bold=False, align=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    for run in p.runs:
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = font_name
    return box

def add_shape(slide, shape_type, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

# ===== 第1页：封面 =====
def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    
    # 装饰大圆
    add_shape(slide, MSO_SHAPE.OVAL, -1, 1, 4, 4, fill_color=RGBColor(0xFF,0x45,0x00), line_color=None)
    add_shape(slide, MSO_SHAPE.OVAL, 7.5, 4, 3, 3, fill_color=RGBColor(0xFF,0x45,0x00), line_color=None)
    
    # 主标题
    add_text(slide, 1, 2.2, 8, 1.5, "群像·星火", 72, C_ORANGE, True, PP_ALIGN.CENTER)
    # 副标题
    add_text(slide, 1, 3.5, 8, 0.8, "基于真实职业经验的多人协同创作平台", 28, C_LIGHT, False, PP_ALIGN.CENTER)
    # Slogan
    add_text(slide, 1, 4.4, 8, 0.6, "让真实的人，在真实的情境中，碰撞出真实的火花", 18, C_GOLD, False, PP_ALIGN.CENTER)
    # 装饰线
    line = add_shape(slide, MSO_SHAPE.RECTANGLE, 3.5, 5.3, 3, 0.04, fill_color=C_ORANGE)
    line.line.fill.background()
    # 底部标签
    add_text(slide, 1, 6.2, 8, 0.5, "TDD v4.0  |  216 Tests Passed  |  Next.js 16 + DeepSeek AI", 14, C_GRAY, False, PP_ALIGN.CENTER)

# ===== 第2页：痛点 =====
def slide_pain(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    
    # 左侧大问号形状
    q_shape = add_shape(slide, MSO_SHAPE.OVAL, 0.8, 1.8, 3.5, 3.5, fill_color=RGBColor(0xFF,0x45,0x00))
    q_shape.line.fill.background()
    add_text(slide, 0.8, 2.8, 3.5, 1, "?", 120, C_DARK, True, PP_ALIGN.CENTER)
    
    # 右侧标题
    add_text(slide, 5, 1.2, 4.5, 0.8, "创作者的困境", 40, C_ORANGE, True)
    
    # 三个痛点卡片
    pains = [
        ("专业细节", "写不对", "急诊流程 / 律师质证 / 外卖跑单"),
        ("视角单一", "写不深", "编剧只采访了医生，没采访护士"),
        ("有经历的人", "没渠道", "退休阿姨 / 急诊护士 / 程序员"),
    ]
    for i, (title, sub, desc) in enumerate(pains):
        y = 2.2 + i * 1.5
        # 卡片背景
        card = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 5, y, 4.5, 1.2, fill_color=C_CARD)
        card.line.color.rgb = RGBColor(0x33,0x33,0x44)
        # 标题
        add_text(slide, 5.3, y+0.15, 2, 0.5, f"{title}，{sub}", 22, C_ORANGE, True)
        # 描述
        add_text(slide, 5.3, y+0.6, 4, 0.5, desc, 14, C_GRAY)

# ===== 第3页：产品定位 =====
def slide_product(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    
    add_text(slide, 1, 0.8, 8, 0.8, "群像·星火", 48, C_ORANGE, True, PP_ALIGN.CENTER)
    add_text(slide, 1, 1.5, 8, 0.5, "基于真实职业经验的多人协同创作平台", 20, C_GRAY, False, PP_ALIGN.CENTER)
    
    # 中央概念图：两个圆形碰撞
    # 左圆
    left = add_shape(slide, MSO_SHAPE.OVAL, 2.5, 2.8, 2, 2, fill_color=RGBColor(0x3B,0x82,0xF6))
    left.line.fill.background()
    add_text(slide, 2.5, 3.5, 2, 0.5, "医生", 20, C_LIGHT, True, PP_ALIGN.CENTER)
    # 右圆
    right = add_shape(slide, MSO_SHAPE.OVAL, 5.5, 2.8, 2, 2, fill_color=RGBColor(0xEC,0x48,0x99))
    right.line.fill.background()
    add_text(slide, 5.5, 3.5, 2, 0.5, "导演", 20, C_LIGHT, True, PP_ALIGN.CENTER)
    # 碰撞火花
    spark = add_shape(slide, MSO_SHAPE.OVAL, 4.3, 3.3, 1.4, 1.4, fill_color=C_ORANGE)
    spark.line.fill.background()
    add_text(slide, 4.3, 3.7, 1.4, 0.5, "碰撞", 16, C_DARK, True, PP_ALIGN.CENTER)
    
    # 底部关键词
    keywords = ["真实身份", "同一情境", "即时对白", "火花标记", "AI串联"]
    for i, kw in enumerate(keywords):
        x = 1.2 + i * 1.6
        tag = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 5.8, 1.4, 0.5, fill_color=C_CARD)
        tag.line.color.rgb = C_ORANGE
        add_text(slide, x, 5.88, 1.4, 0.4, kw, 14, C_ORANGE, False, PP_ALIGN.CENTER)

# ===== 第4页：单人模式 =====
def slide_single(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    
    add_text(slide, 1, 0.6, 8, 0.6, "模式一：单人灵感积累", 36, C_ORANGE, True, PP_ALIGN.CENTER)
    
    # 手机模拟界面
    phone = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 3.5, 1.5, 3, 5, fill_color=RGBColor(0x11,0x11,0x22))
    phone.line.color.rgb = RGBColor(0x44,0x44,0x55)
    phone.line.width = Pt(3)
    
    # 手机内内容
    add_text(slide, 3.7, 1.8, 2.6, 0.4, "深夜急诊室", 14, C_LIGHT, True, PP_ALIGN.CENTER)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 3.8, 2.3, 2.4, 1.2, fill_color=C_CARD)
    add_text(slide, 3.9, 2.5, 2.2, 0.8, "外卖员因\n过度劳累晕倒", 12, C_GRAY, False, PP_ALIGN.CENTER)
    
    # AI 提示气泡
    bubble = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 3.8, 3.8, 2.4, 0.8, fill_color=RGBColor(0xFF,0x45,0x00))
    bubble.line.fill.background()
    add_text(slide, 3.9, 3.95, 2.2, 0.5, "你首先会关注\n哪些生命体征？", 11, C_DARK, False, PP_ALIGN.CENTER)
    
    # 语音按钮模拟
    mic = add_shape(slide, MSO_SHAPE.OVAL, 4.7, 4.9, 0.6, 0.6, fill_color=RGBColor(0x33,0x33,0x44))
    mic.line.color.rgb = C_ORANGE
    add_text(slide, 4.7, 5.0, 0.6, 0.4, "mic", 10, C_ORANGE, False, PP_ALIGN.CENTER)
    
    # 左侧步骤
    steps = [("1", "选身份"), ("2", "刷脑洞"), ("3", "AI催化"), ("4", "存素材")]
    for i, (num, label) in enumerate(steps):
        y = 1.8 + i * 1.2
        circle = add_shape(slide, MSO_SHAPE.OVAL, 0.8, y, 0.6, 0.6, fill_color=C_ORANGE)
        circle.line.fill.background()
        add_text(slide, 0.8, y+0.1, 0.6, 0.4, num, 16, C_DARK, True, PP_ALIGN.CENTER)
        add_text(slide, 1.6, y+0.1, 1.5, 0.4, label, 18, C_LIGHT)
    
    # 右侧数据
    add_text(slide, 7, 2, 2.5, 0.8, "5 类", 48, C_ORANGE, True, PP_ALIGN.RIGHT)
    add_text(slide, 7, 2.7, 2.5, 0.4, "职业标签", 14, C_GRAY, False, PP_ALIGN.RIGHT)
    add_text(slide, 7, 3.5, 2.5, 0.8, "60+", 48, C_ORANGE, True, PP_ALIGN.RIGHT)
    add_text(slide, 7, 4.2, 2.5, 0.4, "引导问题", 14, C_GRAY, False, PP_ALIGN.RIGHT)

# ===== 第5页：双人模式 =====
def slide_duo(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    
    add_text(slide, 1, 0.6, 8, 0.6, "模式二：双人即兴碰撞", 36, C_ORANGE, True, PP_ALIGN.CENTER)
    
    # 左用户
    left = add_shape(slide, MSO_SHAPE.OVAL, 1.5, 2.2, 1.8, 1.8, fill_color=RGBColor(0x3B,0x82,0xF6))
    left.line.fill.background()
    add_text(slide, 1.5, 2.9, 1.8, 0.5, "医生", 20, C_LIGHT, True, PP_ALIGN.CENTER)
    add_text(slide, 1.2, 4.2, 2.4, 0.6, "先推肾上腺素！\n患者室颤了", 14, C_LIGHT, False, PP_ALIGN.CENTER)
    
    # 右用户
    right = add_shape(slide, MSO_SHAPE.OVAL, 6.7, 2.2, 1.8, 1.8, fill_color=RGBColor(0xEC,0x48,0x99))
    right.line.fill.background()
    add_text(slide, 6.7, 2.9, 1.8, 0.5, "导演", 20, C_LIGHT, True, PP_ALIGN.CENTER)
    add_text(slide, 6.8, 4.2, 2.4, 0.6, "等等，情绪高潮\n应该在后面", 14, C_LIGHT, False, PP_ALIGN.CENTER)
    
    # 中间连接 + 火花
    add_shape(slide, MSO_SHAPE.RECTANGLE, 3.5, 3.1, 3, 0.08, fill_color=C_GRAY)
    spark = add_shape(slide, MSO_SHAPE.OVAL, 4.6, 2.8, 1, 1, fill_color=C_ORANGE)
    spark.line.fill.background()
    add_text(slide, 4.6, 3.15, 1, 0.4, "碰撞", 14, C_DARK, True, PP_ALIGN.CENTER)
    
    # 火花标记
    star = add_shape(slide, MSO_SHAPE.STAR_5_POINT, 4.8, 4.2, 0.6, 0.6, fill_color=C_GOLD)
    star.line.fill.background()
    add_text(slide, 4.5, 4.9, 1.2, 0.4, "标记为火花", 12, C_GOLD, False, PP_ALIGN.CENTER)
    
    # 底部数据
    add_text(slide, 1, 5.8, 2.5, 0.6, "60秒", 36, C_ORANGE, True)
    add_text(slide, 1, 6.3, 2.5, 0.3, "匹配等待", 12, C_GRAY)
    add_text(slide, 3.8, 5.8, 2.5, 0.6, "WebSocket", 36, C_ORANGE, True, PP_ALIGN.CENTER)
    add_text(slide, 3.8, 6.3, 2.5, 0.3, "毫秒级同步", 12, C_GRAY, False, PP_ALIGN.CENTER)
    add_text(slide, 6.5, 5.8, 2.5, 0.6, "AI串联", 36, C_ORANGE, True, PP_ALIGN.RIGHT)
    add_text(slide, 6.5, 6.3, 2.5, 0.3, "剧本自动生成", 12, C_GRAY, False, PP_ALIGN.RIGHT)

# ===== 第6页：多人模式 =====
def slide_multi(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    
    add_text(slide, 1, 0.6, 8, 0.6, "模式三：多人剧本共创", 36, C_ORANGE, True, PP_ALIGN.CENTER)
    
    # 导演（大）
    director = add_shape(slide, MSO_SHAPE.OVAL, 4.1, 1.5, 1.8, 1.8, fill_color=C_ORANGE)
    director.line.fill.background()
    add_text(slide, 4.1, 2.2, 1.8, 0.5, "导演", 20, C_DARK, True, PP_ALIGN.CENTER)
    
    # 参与者（围绕）
    roles = [(2.2, 2.5, "医生"), (6.8, 2.5, "护士"), (1.8, 4.2, "家属"), (7.2, 4.2, "律师")]
    for x, y, role in roles:
        p = add_shape(slide, MSO_SHAPE.OVAL, x, y, 1.2, 1.2, fill_color=C_CARD)
        p.line.color.rgb = RGBColor(0x44,0x44,0x55)
        add_text(slide, x, y+0.4, 1.2, 0.4, role, 14, C_LIGHT, False, PP_ALIGN.CENTER)
    
    # 连接线
    for x, y, _ in roles:
        cx, cy = 5.0, 2.4
        px, py = x+0.6, y+0.6
        # 简化：用虚线效果的小矩形代替
        pass
    
    # 导演控场按钮
    controls = [("暂停", 3.2), ("投票", 4.6), ("杀青", 6.0)]
    for label, x in controls:
        btn = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 4.8, 1.2, 0.5, fill_color=C_CARD)
        btn.line.color.rgb = C_ORANGE
        add_text(slide, x, 4.88, 1.2, 0.4, label, 14, C_ORANGE, False, PP_ALIGN.CENTER)
    
    # 底部流程
    flow = ["创建副本", "认领角色", "导演控场", "投票决策", "AI串联", "署名墙"]
    for i, step in enumerate(flow):
        x = 0.6 + i * 1.5
        box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 6.0, 1.3, 0.45, fill_color=C_CARD)
        box.line.color.rgb = RGBColor(0x44,0x44,0x55)
        add_text(slide, x, 6.08, 1.3, 0.35, step, 11, C_LIGHT, False, PP_ALIGN.CENTER)

# ===== 第7页：技术架构 =====
def slide_tech(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    
    add_text(slide, 1, 0.6, 8, 0.6, "技术架构", 36, C_ORANGE, True, PP_ALIGN.CENTER)
    
    # 分层架构图
    layers = [
        ("Next.js 16 + React 19 + Tailwind v4", RGBColor(0x3B,0x82,0xF6), 0.8),
        ("API Routes + NextAuth + Zod", RGBColor(0x8B,0x5C,0xF6), 0.8),
        ("Socket.io + match-engine + room-manager", RGBColor(0xEC,0x48,0x99), 0.8),
        ("Prisma 7 + SQLite", RGBColor(0x22,0xC5,0x5E), 0.8),
    ]
    for i, (label, color, w) in enumerate(layers):
        y = 1.5 + i * 1.1
        rect = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 2.5, y, 5, 0.9, fill_color=color)
        rect.line.fill.background()
        add_text(slide, 2.5, y+0.25, 5, 0.5, label, 18, C_LIGHT, True, PP_ALIGN.CENTER)
    
    # 左侧标签
    labels = ["前端", "API层", "业务层", "数据层"]
    for i, label in enumerate(labels):
        y = 1.5 + i * 1.1
        add_text(slide, 0.8, y+0.25, 1.5, 0.5, label, 16, C_GRAY, False, PP_ALIGN.RIGHT)
    
    # 右侧 DeepSeek 图标
    ds = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 8, 2, 1.5, 1.5, fill_color=RGBColor(0x11,0x44,0x77))
    ds.line.fill.background()
    add_text(slide, 8, 2.5, 1.5, 0.5, "DeepSeek\nAI", 14, C_LIGHT, True, PP_ALIGN.CENTER)
    
    # 底部大数字
    add_text(slide, 1, 5.8, 3, 1, "216", 72, C_GREEN, True)
    add_text(slide, 1, 6.6, 3, 0.4, "Tests Passed", 16, C_GRAY)
    add_text(slide, 4, 5.8, 3, 1, "23", 72, C_ORANGE, True)
    add_text(slide, 4, 6.6, 3, 0.4, "Test Files", 16, C_GRAY)
    add_text(slide, 7, 5.8, 2, 1, "0", 72, C_GREEN, True)
    add_text(slide, 7, 6.6, 2, 0.4, "Failed", 16, C_GRAY)

# ===== 第8页：商业模式 =====
def slide_biz(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    
    add_text(slide, 1, 0.6, 8, 0.6, "商业模式", 36, C_ORANGE, True, PP_ALIGN.CENTER)
    
    # 金字塔三层
    # 底层最大
    b1 = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.5, 4.5, 7, 1.2, fill_color=RGBColor(0x22,0xC5,0x5E))
    b1.line.fill.background()
    add_text(slide, 1.5, 4.75, 7, 0.8, "C 端免费增值", 28, C_DARK, True, PP_ALIGN.CENTER)
    add_text(slide, 1.5, 5.3, 7, 0.4, "基础免费 + AI高级功能付费", 14, C_DARK, False, PP_ALIGN.CENTER)
    
    # 中层
    b2 = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 2.5, 3.0, 5, 1.2, fill_color=RGBColor(0x3B,0x82,0xF6))
    b2.line.fill.background()
    add_text(slide, 2.5, 3.25, 5, 0.8, "B 端内容采购", 28, C_DARK, True, PP_ALIGN.CENTER)
    add_text(slide, 2.5, 3.8, 5, 0.4, "微短剧公司 / 互动小说平台", 14, C_DARK, False, PP_ALIGN.CENTER)
    
    # 顶层
    b3 = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 3.5, 1.5, 3, 1.2, fill_color=C_ORANGE)
    b3.line.fill.background()
    add_text(slide, 3.5, 1.75, 3, 0.8, "IP 共创分润", 28, C_DARK, True, PP_ALIGN.CENTER)
    add_text(slide, 3.5, 2.3, 3, 0.4, "剧本 / 短剧 / 有声书", 14, C_DARK, False, PP_ALIGN.CENTER)

# ===== 第9页：里程碑 =====
def slide_milestone(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    
    add_text(slide, 1, 0.6, 8, 0.6, "里程碑", 36, C_ORANGE, True, PP_ALIGN.CENTER)
    
    # 时间轴
    add_shape(slide, MSO_SHAPE.RECTANGLE, 1, 3.5, 8, 0.06, fill_color=C_GRAY)
    
    phases = [
        (1.5, "P1", "匹配引擎", "9 tests"),
        (3.0, "P2", "房间管理", "16 tests"),
        (4.5, "P3", "WebSocket\n+ AI串联", "21 tests"),
        (6.5, "P4", "TDD全覆盖\n+ AI催化", "216 tests"),
    ]
    for x, label, title, tests in phases:
        # 节点
        node = add_shape(slide, MSO_SHAPE.OVAL, x, 3.3, 0.4, 0.4, fill_color=C_ORANGE)
        node.line.fill.background()
        # 标签
        add_text(slide, x-0.3, 2.5, 1, 0.5, label, 18, C_ORANGE, True, PP_ALIGN.CENTER)
        add_text(slide, x-0.5, 3.9, 1.4, 0.8, title, 14, C_LIGHT, False, PP_ALIGN.CENTER)
        add_text(slide, x-0.5, 4.5, 1.4, 0.4, tests, 12, C_GREEN, False, PP_ALIGN.CENTER)
    
    # 大数字
    add_text(slide, 1, 5.5, 8, 1, "216 / 216", 64, C_GREEN, True, PP_ALIGN.CENTER)
    add_text(slide, 1, 6.3, 8, 0.4, "All Tests Passed  |  23 Test Files  |  22 API Routes Covered", 14, C_GRAY, False, PP_ALIGN.CENTER)

# ===== 第10页：结尾 =====
def slide_end(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    
    # 装饰圆
    add_shape(slide, MSO_SHAPE.OVAL, -1, -1, 4, 4, fill_color=RGBColor(0xFF,0x45,0x00))
    add_shape(slide, MSO_SHAPE.OVAL, 7, 5, 3, 3, fill_color=RGBColor(0xFF,0x45,0x00))
    
    add_text(slide, 1, 2.0, 8, 1, "最好的故事", 64, C_ORANGE, True, PP_ALIGN.CENTER)
    add_text(slide, 1, 3.0, 8, 0.6, "不是一个人关在房间里写出来的", 24, C_LIGHT, False, PP_ALIGN.CENTER)
    add_text(slide, 1, 3.7, 8, 0.6, "而是让真实的人在真实的情境中碰撞出来的", 20, C_GOLD, False, PP_ALIGN.CENTER)
    add_text(slide, 1, 4.8, 8, 0.8, "谢谢大家", 40, C_ORANGE, True, PP_ALIGN.CENTER)
    add_text(slide, 1, 5.6, 8, 0.4, "github.com/qunxiang-xinghuo", 14, C_GRAY, False, PP_ALIGN.CENTER)

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    slide_cover(prs)
    slide_pain(prs)
    slide_product(prs)
    slide_single(prs)
    slide_duo(prs)
    slide_multi(prs)
    slide_tech(prs)
    slide_biz(prs)
    slide_milestone(prs)
    slide_end(prs)
    
    output_path = "docs/群像星火-路演PPT-v2.pptx"
    prs.save(output_path)
    print(f"[OK] PPT v2 generated: {output_path}")
    print(f"     Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
